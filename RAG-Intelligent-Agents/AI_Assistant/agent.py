import base64
import pandas as pd
import json
import fitz  # PyMuPDF
import nbformat
from docx import Document
from pptx import Presentation
from sympy import sympify
import pint
import os
import ssl
ssl._create_default_https_context = ssl._create_unverified_context

from langchain_core.tools import tool
from langchain_core.messages import HumanMessage, SystemMessage, AnyMessage
from langchain_openai import ChatOpenAI
from langgraph.graph import StateGraph, START
from langgraph.graph.message import add_messages
from langgraph.prebuilt import ToolNode, tools_condition
from typing import TypedDict, Annotated, Optional, List
from langchain_community.utilities import WikipediaAPIWrapper, SerpAPIWrapper
from langchain_community.tools import WikipediaQueryRun
from langchain_community.tools.riza.command import ExecPython
from llama_index.core.tools import FunctionTool
import whisper

# Import API keys from secret_key.py
from secret_key import RIZA_API_KEY, OPENAI_API_KEY, SERPAPI_API_KEY
# Import token utilities for safe processing
from token_utils import estimate_tokens, truncate_text_smart, get_safe_prompt

# Set environment variables
os.environ["RIZA_API_KEY"] = RIZA_API_KEY
os.environ["OPENAI_API_KEY"] = OPENAI_API_KEY

# Correctly initialize with the named parameter
serp_wrapper     = SerpAPIWrapper(serpapi_api_key=SERPAPI_API_KEY)
# Initialize the wrapper (no API key needed)
wiki_api = WikipediaAPIWrapper()

# Create the tool
wiki_tool = WikipediaQueryRun(api_wrapper=wiki_api)

# Load Whisper model at startup to avoid delays later
print("🔊 Loading Whisper model...")
try:
    whisper_model = whisper.load_model("base")
    print("✅ Whisper model loaded successfully!")
except Exception as e:
    print(f"⚠️ Warning: Could not load Whisper model: {e}")
    whisper_model = None

# Define State
MessagesType = Annotated[List[AnyMessage], add_messages]

class AgentState(TypedDict):
    input_file: Optional[str]
    messages: MessagesType

# Vision LLM for OCR
vision_llm = ChatOpenAI(model="gpt-4o")

@tool
def extract_text(img_path: str) -> str:
    """
    Extract text from an image using a vision-capable model.
    Input: local image file path (e.g. /content/1.jpg)
    Output: Extracted text from the image (no explanations).
    """
    try:
        with open(img_path, "rb") as image_file:
            image_bytes = image_file.read()
        image_base64 = base64.b64encode(image_bytes).decode("utf-8")
        message = [
            HumanMessage(content=[
                {
                    "type": "text",
                    "text": "Extract all the text from this image. Return only the extracted text, no explanations."
                },
                {
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/png;base64,{image_base64}"
                    }
                }
            ])
        ]
        response = vision_llm.invoke(message)
        return response.content.strip()
    except Exception as e:
        return f"Error extracting text: {str(e)}"

@tool
def search_query(query: str) -> str:
    """Search for real-time data using SerpAPI"""
    return serp_wrapper.run(query)

@tool
def read_file(file_path: str) -> str:
    """
    Reads and extracts text from supported file types (.txt, .csv, .json, .pdf, .docx, .pptx, .xlsx, .ipynb, .mp3).
    Does NOT handle images (use extract_text for images). Returns readable string or error message.
    Automatically handles large files by summarizing or chunking content.
    """
    try:
        content = ""
        
        if file_path.endswith(".txt"):
            with open(file_path, "r") as f:
                content = f.read()

        elif file_path.endswith(".csv"):
            df = pd.read_csv(file_path)
            # For large CSV files, provide summary instead of full content
            if len(df) > 100:  # More aggressive chunking - 100 rows instead of 1000
                content = f"CSV File Summary (Large file with {len(df)} rows):\n"
                content += f"- Rows: {len(df)}\n"
                content += f"- Columns: {len(df.columns)}\n"
                content += f"- Column names: {list(df.columns)}\n"
                content += f"- Data types:\n{df.dtypes.to_string()}\n"
                content += f"- Missing values:\n{df.isnull().sum().to_string()}\n"
                content += f"- First 5 rows:\n{df.head(5).to_string()}\n"
                content += f"- Last 5 rows:\n{df.tail(5).to_string()}\n"
                
                # Add basic statistics but limit to numeric columns only
                numeric_cols = df.select_dtypes(include=['number']).columns
                if len(numeric_cols) > 0:
                    content += f"- Basic statistics (numeric columns only):\n{df[numeric_cols].describe().to_string()}\n"
                
                # Add sample of unique values for categorical columns
                categorical_cols = df.select_dtypes(include=['object']).columns
                if len(categorical_cols) > 0:
                    content += f"- Sample categorical values:\n"
                    for col in categorical_cols[:5]:  # Limit to first 5 categorical columns
                        unique_vals = df[col].unique()[:10]  # First 10 unique values
                        content += f"  {col}: {list(unique_vals)}\n"
            else:
                content = df.to_string()

        elif file_path.endswith(".json"):
            with open(file_path, "r") as f:
                data = json.load(f)
                json_str = json.dumps(data, indent=2)
                # If JSON is too large, provide structure summary
                if len(json_str) > 10000:
                    content = f"Large JSON File Summary:\n"
                    content += f"- Total characters: {len(json_str)}\n"
                    content += f"- Root type: {type(data).__name__}\n"
                    if isinstance(data, dict):
                        content += f"- Keys: {list(data.keys())[:20]}...\n"  # First 20 keys
                    elif isinstance(data, list):
                        content += f"- Array length: {len(data)}\n"
                        content += f"- First few items: {str(data[:5])}\n"
                    content += f"- Sample content (first 2000 chars):\n{json_str[:2000]}...\n"
                else:
                    content = json_str

        elif file_path.endswith(".pdf"):
            doc = fitz.open(file_path)
            full_text = "\n".join(page.get_text() for page in doc)
            # If PDF is too large, provide summary
            if len(full_text) > 15000:
                content = f"Large PDF Summary:\n"
                content += f"- Total pages: {len(doc)}\n"
                content += f"- Total characters: {len(full_text)}\n"
                content += f"- First 5000 characters:\n{full_text[:5000]}...\n"
                content += f"- Last 2000 characters:\n...{full_text[-2000:]}\n"
            else:
                content = full_text

        elif file_path.endswith(".docx"):
            doc = Document(file_path)
            full_text = "\n".join(p.text for p in doc.paragraphs)
            # Handle large documents
            if len(full_text) > 15000:
                content = f"Large Document Summary:\n"
                content += f"- Total paragraphs: {len(doc.paragraphs)}\n"
                content += f"- Total characters: {len(full_text)}\n"
                content += f"- First 5000 characters:\n{full_text[:5000]}...\n"
                content += f"- Last 2000 characters:\n...{full_text[-2000:]}\n"
            else:
                content = full_text

        elif file_path.endswith(".pptx"):
            prs = Presentation(file_path)
            slides_text = []
            for i, slide in enumerate(prs.slides):
                slide_text = f"Slide {i+1}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                slides_text.append(slide_text)
            
            full_text = "\n".join(slides_text)
            # Handle large presentations
            if len(full_text) > 15000:
                content = f"Large Presentation Summary:\n"
                content += f"- Total slides: {len(prs.slides)}\n"
                content += f"- Total characters: {len(full_text)}\n"
                content += f"- First 10 slides:\n{chr(10).join(slides_text[:10])}\n"
            else:
                content = full_text

        elif file_path.endswith(".ipynb"):
            with open(file_path) as f:
                nb = nbformat.read(f, as_version=4)
                cells_content = []
                for cell in nb.cells:
                    if cell.cell_type in ['code', 'markdown']:
                        cells_content.append(f"{cell.cell_type}: {cell['source']}")
                
                full_text = "\n".join(cells_content)
                # Handle large notebooks
                if len(full_text) > 15000:
                    content = f"Large Notebook Summary:\n"
                    content += f"- Total cells: {len(nb.cells)}\n"
                    content += f"- Cell types: {[cell.cell_type for cell in nb.cells]}\n"
                    content += f"- First 5000 characters:\n{full_text[:5000]}...\n"
                else:
                    content = full_text

        elif file_path.endswith(".xlsx"):
            df = pd.read_excel(file_path, engine="openpyxl")
            # Handle large Excel files similar to CSV - more aggressive chunking
            if len(df) > 100:  # More aggressive chunking
                content = f"Excel File Summary (Large file with {len(df)} rows):\n"
                content += f"- Rows: {len(df)}\n"
                content += f"- Columns: {len(df.columns)}\n"
                content += f"- Column names: {list(df.columns)}\n"
                content += f"- Data types:\n{df.dtypes.to_string()}\n"
                content += f"- Missing values:\n{df.isnull().sum().to_string()}\n"
                content += f"- First 5 rows:\n{df.head(5).to_string()}\n"
                content += f"- Last 5 rows:\n{df.tail(5).to_string()}\n"
                
                # Add basic statistics but limit to numeric columns only
                numeric_cols = df.select_dtypes(include=['number']).columns
                if len(numeric_cols) > 0:
                    content += f"- Basic statistics (numeric columns only):\n{df[numeric_cols].describe().to_string()}\n"
                
                # Add sample of unique values for categorical columns
                categorical_cols = df.select_dtypes(include=['object']).columns
                if len(categorical_cols) > 0:
                    content += f"- Sample categorical values:\n"
                    for col in categorical_cols[:5]:  # Limit to first 5 categorical columns
                        unique_vals = df[col].unique()[:10]  # First 10 unique values
                        content += f"  {col}: {list(unique_vals)}\n"
            else:
                content = df.to_string()

        elif file_path.endswith(".xls"):
            return "❌ Legacy .xls files are not fully supported. Please re-save as .xlsx."

        elif file_path.endswith((".mp3", ".wav", ".m4a", ".flac")):
            # Use the transcribe_mp3 function for audio files
            print(f"🎵 Processing audio file: {file_path}")
            return transcribe_mp3(file_path)

        else:
            return "⚠️ Unsupported file type."

        # Final check for content length using token estimation
        estimated_tokens = estimate_tokens(content)
        
        if estimated_tokens > 20000:  # More aggressive limit - 20k tokens max
            print(f"⚠️ Large content detected: {estimated_tokens} tokens, truncating...")
            content = truncate_text_smart(content, 15000)  # Truncate to 15k tokens
            print(f"✅ Content truncated to ~{estimate_tokens(content)} tokens")

        return content

    except Exception as e:
        return f"❌ Error while reading file: {e}"


@tool
def solve_math(expression: str) -> str:
    """Solve symbolic math expressions."""
    try:
        return str(sympify(expression))
    except Exception as e:
        return f"Failed to parse math: {e}"

@tool
def convert_units(query: str) -> str:
    """Convert units. Example: '90 km/h to m/s'"""
    ureg = pint.UnitRegistry()
    try:
        if " to " not in query:
            return "Invalid format. Use 'X unit to new_unit' format."
        value_part, target_unit = query.split(" to ")
        value = ureg.Quantity(value_part.strip())
        return str(value.to(target_unit.strip()))
    except Exception as e:
        return f"Conversion failed: {e}"

@tool
def transcribe_mp3(mp3_path: str) -> str:
    """
    Transcribe an audio file to text using Whisper.
    Supports MP3, WAV, M4A, FLAC, and other audio formats.

    Args:
        mp3_path: Path to a local audio file.

    Returns:
        Transcribed text or error message.
    """
    global whisper_model
    try:
        # Check if file exists
        if not os.path.exists(mp3_path):
            return f"❌ Error: Audio file not found at {mp3_path}"
        
        # Check file size (optional limit)
        file_size = os.path.getsize(mp3_path) / (1024 * 1024)  # Size in MB
        print(f"🎵 Processing audio file: {mp3_path} ({file_size:.2f} MB)")
        
        # Load the model on first use
        if whisper_model is None:
            print("🔊 Loading Whisper model (this may take a moment)...")
            whisper_model = whisper.load_model("base")
            print("✅ Whisper model loaded successfully!")
        
        # Transcribe the audio file
        print("🎧 Starting transcription...")
        result = whisper_model.transcribe(
            mp3_path,
            verbose=False,  # Reduce console output
            fp16=False,     # Better compatibility
        )
        
        transcribed_text = result["text"].strip()
        
        if not transcribed_text:
            return "⚠️ Warning: No speech detected in the audio file. The file might be silent or contain only noise."
        
        print(f"✅ Transcription completed: {len(transcribed_text)} characters")
        return transcribed_text
        
    except Exception as e:
        error_msg = f"❌ Error during transcription: {str(e)}"
        print(error_msg)
        return error_msg

@tool
def analyze_large_dataset(file_path: str, analysis_question: str) -> str:
    """
    Specialized tool for analyzing large datasets with smart chunking and summarization.
    
    Args:
        file_path: Path to the dataset file
        analysis_question: What analysis to perform
    
    Returns:
        Analysis results optimized for large files
    """
    try:
        # First, get basic file info
        if not os.path.exists(file_path):
            return f"❌ File not found: {file_path}"
        
        file_size = os.path.getsize(file_path) / (1024 * 1024)  # Size in MB
        print(f"📊 Analyzing large dataset: {file_path} ({file_size:.2f} MB)")
        
        # Load based on file type
        if file_path.endswith('.csv'):
            df = pd.read_csv(file_path)
        elif file_path.endswith('.xlsx'):
            df = pd.read_excel(file_path, engine="openpyxl")
        else:
            return "❌ Unsupported file type for dataset analysis. Use CSV or Excel files."
        
        # Generate focused analysis based on question
        analysis_results = []
        
        # Basic dataset info
        analysis_results.append(f"📋 Dataset Overview:")
        analysis_results.append(f"- Shape: {df.shape[0]} rows × {df.shape[1]} columns")
        analysis_results.append(f"- Columns: {', '.join(df.columns.tolist())}")
        analysis_results.append(f"- Memory usage: {df.memory_usage(deep=True).sum() / 1024 / 1024:.2f} MB")
        
        # Data types and missing values
        analysis_results.append(f"\n📊 Data Quality:")
        analysis_results.append(f"- Data types: {df.dtypes.value_counts().to_dict()}")
        missing_data = df.isnull().sum()
        if missing_data.any():
            analysis_results.append(f"- Missing values: {missing_data[missing_data > 0].to_dict()}")
        else:
            analysis_results.append("- Missing values: None")
        
        # Focus on the specific analysis question
        if "revenue" in analysis_question.lower() or "financial" in analysis_question.lower():
            # Financial analysis
            numeric_cols = df.select_dtypes(include=['number']).columns
            if len(numeric_cols) > 0:
                analysis_results.append(f"\n💰 Financial Analysis:")
                for col in numeric_cols:
                    if any(term in col.lower() for term in ['revenue', 'sales', 'profit', 'cost', 'amount', 'price']):
                        analysis_results.append(f"- {col}: ${df[col].sum():,.2f} total, ${df[col].mean():,.2f} average")
        
        elif "trend" in analysis_question.lower() or "time" in analysis_question.lower():
            # Time series analysis
            date_cols = df.select_dtypes(include=['datetime', 'object']).columns
            potential_date_cols = [col for col in date_cols if any(term in col.lower() for term in ['date', 'time', 'created', 'updated'])]
            
            if potential_date_cols:
                analysis_results.append(f"\n📈 Trend Analysis:")
                for col in potential_date_cols[:2]:  # Limit to first 2 date columns
                    try:
                        df[col] = pd.to_datetime(df[col])
                        analysis_results.append(f"- {col}: {df[col].min()} to {df[col].max()}")
                    except:
                        continue
        
        elif "customer" in analysis_question.lower():
            # Customer analysis
            potential_customer_cols = [col for col in df.columns if any(term in col.lower() for term in ['customer', 'client', 'user', 'account'])]
            if potential_customer_cols:
                analysis_results.append(f"\n👥 Customer Analysis:")
                for col in potential_customer_cols[:3]:  # Limit to first 3 customer columns
                    unique_count = df[col].nunique()
                    analysis_results.append(f"- {col}: {unique_count} unique values")
        
        # Summary statistics for numeric columns
        numeric_cols = df.select_dtypes(include=['number']).columns
        if len(numeric_cols) > 0:
            analysis_results.append(f"\n📊 Key Statistics:")
            for col in numeric_cols[:5]:  # Limit to first 5 numeric columns
                analysis_results.append(f"- {col}: min={df[col].min():.2f}, max={df[col].max():.2f}, avg={df[col].mean():.2f}")
        
        # Top categories for categorical data
        categorical_cols = df.select_dtypes(include=['object']).columns
        if len(categorical_cols) > 0:
            analysis_results.append(f"\n🏷️ Top Categories:")
            for col in categorical_cols[:3]:  # Limit to first 3 categorical columns
                top_values = df[col].value_counts().head(3)
                analysis_results.append(f"- {col}: {top_values.to_dict()}")
        
        # Combine all results
        final_analysis = "\n".join(analysis_results)
        
        # Add specific insights based on the question
        final_analysis += f"\n\n🎯 Specific Insights for: '{analysis_question}'"
        
        return final_analysis
        
    except Exception as e:
        return f"❌ Error analyzing dataset: {str(e)}"

# All tools
tools = [
    extract_text,
    search_query,
    read_file,
    solve_math,
    convert_units,
    transcribe_mp3,
    wiki_tool,
    ExecPython()
]

# LangGraph setup
llm = ChatOpenAI(model="gpt-4o")
llm_with_tools = llm.bind_tools(tools)

def assistant(state: AgentState):
    sys_prompt = (
        "You are a general AI assistant. I will ask you a question. Report your thoughts, and finish your answer with the following template: FINAL ANSWER: [YOUR FINAL ANSWER]. YOUR FINAL ANSWER should be a number OR as few words as possible OR a comma separated list of numbers and/or strings. If you are asked for a number, don't use comma to write your number neither use units such as $ or percent sign unless specified otherwise. If you are asked for a string, don't use articles, neither abbreviations (e.g. for cities), and write the digits in plain text unless specified otherwise. If you are asked for a comma separated list, apply the above rules depending of whether the element to be put in the list is a number or a string."
    )
    file = state.get("input_file", None)
    
    # Create system message with token-aware content
    system_content = sys_prompt + (f" Current input file: {file}" if file else "")
    
    # Check if we need to limit the conversation due to token constraints
    messages = state["messages"]
    total_message_content = " ".join([str(msg.content) for msg in messages])
    total_tokens = estimate_tokens(system_content + total_message_content)
    
    # If approaching token limit, provide a warning
    if total_tokens > 25000:
        print(f"⚠️ High token usage detected: ~{total_tokens} tokens. Using conservative processing.")
        # Optionally truncate older messages or provide summary
        if len(messages) > 3:
            messages = messages[-3:]  # Keep only last 3 messages
            print("📝 Truncated conversation history to manage token limits.")
    
    system = SystemMessage(content=system_content)
    
    try:
        response = llm_with_tools.invoke([system] + messages)
        return {"messages": [response], "input_file": file}
    except Exception as e:
        error_msg = str(e)
        if "rate limit" in error_msg.lower() or "429" in error_msg:
            # Rate limit error - provide helpful message
            error_response = HumanMessage(content=f"❌ Rate limit exceeded. The file or request is too large. Please try with a smaller file or break down your request into smaller parts. Error: {error_msg}")
            return {"messages": [error_response], "input_file": file}
        else:
            # Other errors
            error_response = HumanMessage(content=f"❌ Error processing request: {error_msg}")
            return {"messages": [error_response], "input_file": file}

# Build graph
builder = StateGraph(AgentState)
builder.add_node("assistant", assistant)
builder.add_node("tools", ToolNode(tools))
builder.add_edge(START, "assistant")
builder.add_conditional_edges("assistant", tools_condition)
builder.add_edge("tools", "assistant")
react_graph = builder.compile()
